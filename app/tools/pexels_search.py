"""Pexels 图片搜索工具 - 通过 Pexels API 搜索图片并插入 PPT 幻灯片"""

import json
import logging
import uuid
from pathlib import Path

import httpx
from langchain_core.tools import tool
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from app.core.config import settings
from app.tools.ppt_generator import (
    DEFAULT_THEME,
    SLIDE_WIDTH,
    STYLE_THEMES,
)

logger = logging.getLogger(__name__)

# Pexels 常量
PEXELS_API_URL = "https://api.pexels.com/v1/search"
PEXELS_PER_PAGE = 1
PEXELS_ORIENTATION = "landscape"
PICSUM_URL_TEMPLATE = "https://picsum.photos/800/600?random={}"


def _resolve_ppt_path(filename: str) -> Path:
    return settings.workspace_path / "ppt_output" / filename


def _get_theme(style: str) -> dict:
    return STYLE_THEMES.get(style, DEFAULT_THEME)


async def _search_pexels(keywords: str) -> str | None:
    """调用 Pexels API 搜索图片，返回图片 URL"""
    api_key = settings.pexels_api_key
    if not api_key:
        logger.warning("pexels_api_key 未配置，跳过 Pexels 搜索")
        return None

    url = (
        f"{PEXELS_API_URL}"
        f"?query={keywords}"
        f"&per_page={PEXELS_PER_PAGE}"
        f"&orientation={PEXELS_ORIENTATION}"
    )

    async with httpx.AsyncClient(timeout=30.0) as client:
        try:
            resp = await client.get(url, headers={"Authorization": api_key})
            if resp.status_code != 200:
                logger.error("Pexels API 失败: %d", resp.status_code)
                return None

            photos = resp.json().get("photos", [])
            if not photos:
                logger.warning("Pexels 无结果: %s", keywords)
                return None

            return photos[0].get("src", {}).get("large")
        except Exception as e:
            logger.error("Pexels API 异常: %s", e)
            return None


async def _download_image(url: str, save_dir: Path) -> Path | None:
    """下载图片到本地"""
    save_dir.mkdir(parents=True, exist_ok=True)
    img_path = save_dir / f"img_{uuid.uuid4().hex[:8]}.jpg"

    async with httpx.AsyncClient(timeout=30.0) as client:
        try:
            resp = await client.get(url, follow_redirects=True)
            if resp.status_code == 200:
                img_path.write_bytes(resp.content)
                logger.info(
                    "图片下载成功: %s (%d KB)", img_path.name, len(resp.content) // 1024
                )
                return img_path
        except Exception as e:
            logger.error("图片下载失败: %s", e)
    return None


def _add_image_slide(prs: Presentation, image_path: Path, title: str, theme: dict):
    """向 PPT 添加一张图文幻灯片（标题栏 + 居中大图）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 白色背景
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 顶部标题栏
    header_h = Inches(1.1)
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        header_h,
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme["primary"]
    header.line.fill.background()

    if title:
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(11), header_h - Inches(0.3)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.color.rgb = theme["light_text"]
        p.font.bold = True

    # 图片居中插入
    img_left = Inches(1.5)
    img_top = Inches(1.4)
    img_width = Inches(10.3)
    img_height = Inches(5.5)
    slide.shapes.add_picture(str(image_path), img_left, img_top, img_width, img_height)


async def _prepare_image_asset(keywords: str) -> tuple[Path | None, str]:
    """搜索并下载图片，但不打开或修改任何 PPT 文件。"""
    images_dir = settings.workspace_path / "ppt_output" / "images"
    image_url = await _search_pexels(keywords)
    image_source = "pexels"

    if not image_url:
        fallback_pos = abs(hash(keywords)) % 1000
        image_url = PICSUM_URL_TEMPLATE.format(fallback_pos)
        image_source = "picsum"
        logger.info("降级使用 picsum 图片: %s", image_url)

    image_path = await _download_image(image_url, images_dir)

    if image_source == "picsum" and not image_path:
        retry_url = PICSUM_URL_TEMPLATE.format(fallback_pos + 1)
        logger.info("重试 picsum: %s", retry_url)
        image_path = await _download_image(retry_url, images_dir)

    return image_path, image_source


@tool
async def prepare_image_operation(
    filename: str,
    keywords: str,
    slide_title: str = "",
    style: str = "business",
) -> str:
    """准备一项 PPT 配图操作，但不直接修改 PPT 文件。

    返回的 operation 会由 Workflow 的确定性 ppt_writer_node 汇总并统一写入。

    Args:
        filename: 目标 PPT 文件名，仅用于标识操作目标
        keywords: 图片搜索关键词
        slide_title: 新增图片页标题
        style: PPT 主题风格
    """
    try:
        file_path = _resolve_ppt_path(filename)
        if not file_path.exists():
            return json.dumps(
                {"success": False, "message": f"PPT 文件不存在: {filename}"},
                ensure_ascii=False,
            )

        image_path, image_source = await _prepare_image_asset(keywords)
        if not image_path:
            return json.dumps(
                {"success": False, "message": "图片下载失败，请检查网络连接"},
                ensure_ascii=False,
            )

        operation = {
            "operation_id": f"image-{uuid.uuid4().hex}",
            "type": "image",
            "filename": filename,
            "asset_path": str(image_path),
            "slide_title": slide_title,
            "style": style,
            "keywords": keywords,
            "image_source": image_source,
        }
        return json.dumps(
            {
                "success": True,
                "operation": operation,
                "message": f"配图资源已准备，等待统一写入: {keywords}",
            },
            ensure_ascii=False,
        )
    except Exception as exc:
        logger.exception("准备配图操作失败")
        return json.dumps(
            {"success": False, "message": f"准备配图操作失败: {exc}"},
            ensure_ascii=False,
        )


@tool
async def add_image_slide(
    filename: str,
    keywords: str,
    slide_title: str = "",
    style: str = "business",
) -> str:
    """通过关键词搜索图片并插入为新的PPT幻灯片。

    使用 Pexels API 搜索匹配关键词的高清图片，下载后插入到 PPT 中。
    如果 Pexels 搜索失败或 API Key 未配置，自动降级使用 picsum 随机图片。

    Args:
        filename: PPT文件名（在ppt_output目录下）
        keywords: 图片搜索关键词（英文效果更好），如 "technology", "business meeting", "nature landscape"
        slide_title: 幻灯片标题（可选，不传则无标题栏文字）
        style: 主题风格，可选 business/creative/academic/minimalist
    """
    try:
        theme = _get_theme(style)
        file_path = _resolve_ppt_path(filename)
        if not file_path.exists():
            return f"错误：PPT 文件不存在: {filename}"

        image_path, image_source = await _prepare_image_asset(keywords)

        if not image_path:
            return json.dumps(
                {"success": False, "message": "图片下载失败，请检查网络连接"},
                ensure_ascii=False,
            )

        # 4. 插入幻灯片
        prs = Presentation(str(file_path))
        _add_image_slide(prs, image_path, slide_title, theme)
        prs.save(str(file_path))

        result = {
            "success": True,
            "keywords": keywords,
            "image_source": image_source,
            "message": f"配图幻灯片已添加: {filename}（关键词: {keywords}）",
        }
        logger.info("配图幻灯片添加成功: %s, 关键词: %s", filename, keywords)
        return json.dumps(result, ensure_ascii=False)

    except Exception as e:
        logger.error("配图幻灯片添加失败: %s", e)
        return f"错误：配图幻灯片添加失败: {e}"
