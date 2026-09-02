"""确定性 PPT Writer：汇总并行 Agent 的操作并单点写入文件。"""

import asyncio
import logging
import os
import uuid
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.agents.workflow.state import WorkflowState
from app.core.config import settings
from app.tools.chart_generator import _add_chart_slide_to_ppt
from app.tools.pexels_search import _add_image_slide
from app.tools.ppt_generator import DEFAULT_THEME, STYLE_THEMES

logger = logging.getLogger(__name__)

PPT_WRITER_NODE = "ppt_writer_node"
ASSETS_SKIP_NODE = "assets_skip_node"
DEBUG_ASSETS_JOIN_NODE = "debug_assets_join_node"


def _safe_path(path_value: str, *, root: Path) -> Path:
    """只允许 Workflow 操作 workspace 内的文件，阻止路径穿越。"""
    resolved_root = root.resolve()
    candidate = Path(path_value)
    if not candidate.is_absolute():
        candidate = resolved_root / candidate
    candidate = candidate.resolve()
    if not candidate.is_relative_to(resolved_root):
        raise ValueError(f"资源路径不在 workspace 内: {path_value}")
    return candidate


def _ppt_path(filename: str) -> Path:
    output_dir = (settings.workspace_path / "ppt_output").resolve()
    return _safe_path(filename, root=output_dir)


def _operation_sort_key(operation: dict[str, Any]) -> tuple:
    """并行分支返回顺序不稳定，因此提交前必须确定性排序。"""
    return (
        int(operation.get("slide_index", 10**9)),
        0 if operation.get("type") == "image" else 1,
        str(operation.get("operation_id", "")),
    )


def _manifest_entry_from_operation(operation: dict[str, Any]) -> dict[str, Any]:
    """把已实际写入的资源页转换成可持久化的页面清单项。"""
    operation_type = str(operation.get("type") or "")
    entry: dict[str, Any] = {
        "title": str(operation.get("slide_title") or ""),
        "bullets": [],
        "layout_hint": operation_type,
        "speaker_notes": "",
        "asset_type": operation_type,
    }
    if operation_type == "image":
        entry["keywords"] = str(operation.get("keywords") or "")
        entry["image_source"] = str(operation.get("image_source") or "")
    elif operation_type == "chart":
        entry["chart_type"] = str(operation.get("chart_type") or "")
        entry["chart_title"] = str(operation.get("chart_title") or "")
    return entry


def _apply_operation(prs: Presentation, operation: dict[str, Any]) -> None:
    operation_type = operation.get("type")
    asset_path_value = operation.get("asset_path")
    if not isinstance(asset_path_value, str) or not asset_path_value:
        raise ValueError("PPT 写入操作缺少 asset_path")

    asset_path = _safe_path(asset_path_value, root=settings.workspace_path)
    if not asset_path.is_file():
        raise FileNotFoundError(f"编辑资源不存在: {asset_path.name}")

    style = str(operation.get("style") or "business")
    theme = STYLE_THEMES.get(style, DEFAULT_THEME)
    slide_title = str(operation.get("slide_title") or "")

    if operation_type == "image":
        _add_image_slide(prs, asset_path, slide_title, theme)
        return
    if operation_type == "chart":
        _add_chart_slide_to_ppt(prs, asset_path, slide_title, theme)
        return
    raise ValueError(f"不支持的 PPT 写入操作类型: {operation_type!r}")


def _apply_operations_atomically(
    file_path: Path,
    operations: list[dict[str, Any]],
) -> None:
    """在内存中应用全部操作，通过临时文件 + os.replace 原子提交。"""
    if not file_path.is_file():
        raise FileNotFoundError(f"PPT 文件不存在: {file_path.name}")

    temp_path = file_path.with_name(
        f".{file_path.stem}.{uuid.uuid4().hex}.tmp.pptx"
    )
    try:
        prs = Presentation(str(file_path))
        for operation in operations:
            _apply_operation(prs, operation)
        prs.save(str(temp_path))

        # 提交前重新打开，至少保证生成物是 python-pptx 可解析的 PPTX。
        Presentation(str(temp_path))
        os.replace(temp_path, file_path)
    finally:
        temp_path.unlink(missing_ok=True)


async def ppt_writer_node(state: WorkflowState) -> dict[str, Any]:
    """合并 Image/Chart 操作；本节点不调用 Agent 或 LLM。"""
    filename = state.get("filename")
    if not filename:
        raise ValueError("ppt_writer_node 缺少目标 PPT 文件名")

    operations = sorted(
        state.get("asset_operations", []),
        key=_operation_sort_key,
    )
    if state.get("intent") == "edit":
        requested_types = set(state.get("asset_tasks", []))
        prepared_types = {
            str(operation.get("type"))
            for operation in operations
            if operation.get("type") in {"image", "chart"}
        }
        missing_types = requested_types - prepared_types
        if missing_types:
            missing_text = "、".join(sorted(missing_types))
            return {
                "asset_apply_status": "failed",
                "applied_operation_ids": [],
                "workflow_error": f"Assets 未准备出所需资源：{missing_text}",
            }
    if not operations:
        if state.get("intent") == "edit":
            return {
                "asset_apply_status": "failed",
                "applied_operation_ids": [],
                "workflow_error": "Assets 未产生任何可执行的图片或图表操作",
            }
        return {
            "asset_apply_status": "skipped",
            "applied_operation_ids": [],
            "completed_agents": ["writer"],
            "completed_stages": ["assets"],
        }

    # 当前生产子图中 ppt_writer_node 是唯一 Writer，因此暂不加锁。未来若允许
    # 多个 Workflow 跨 Worker 修改同一 ppt_id，可在这一层接入 Redis Lock。
    file_path = _ppt_path(filename)
    try:
        await asyncio.to_thread(
            _apply_operations_atomically,
            file_path,
            operations,
        )
    except Exception as writer_exc:
        logger.exception(
            "ppt_writer_node 写入增强资源失败，检查基础 PPT 是否可降级交付: %s",
            filename,
        )
        try:
            # 原子写入失败不应破坏旧文件；只有旧文件仍可解析时才允许把
            # Assets/Beautify 从本轮要求中移除并降级为基础版成功交付。
            await asyncio.to_thread(Presentation, str(file_path))
        except Exception as base_ppt_exc:
            logger.exception("基础 PPT 已不可用，无法执行降级交付: %s", filename)
            raise writer_exc from base_ppt_exc

        if state.get("intent") == "edit":
            return {
                "asset_apply_status": "failed",
                "applied_operation_ids": [],
                "workflow_error": f"资源写入失败：{writer_exc}",
            }

        return {
            "asset_apply_status": "failed",
            "applied_operation_ids": [],
            "asset_tasks": [],
            "completed_agents": ["writer"],
            # 阶段已被确定性处理并选择降级，记录为完成可避免断点恢复时
            # 再次规划和执行同一批增强操作。
            "completed_stages": ["assets"],
            "required_stages": [
                stage
                for stage in state.get("required_stages", [])
                if stage not in {"assets", "beautify"}
            ],
            "attempt_error": None,
            "route_reason": "资源写入失败，已降级交付基础版本",
        }

    operation_ids = [
        str(operation["operation_id"])
        for operation in operations
        if operation.get("operation_id")
    ]
    logger.info(
        "ppt_writer_node 已统一写入 PPT: filename=%s operations=%d",
        filename,
        len(operations),
    )
    return {
        "asset_apply_status": "succeeded",
        "applied_operation_ids": operation_ids,
        # 当前底层工具会把资源页追加到文件末尾，因此按同一确定性顺序
        # 追加清单项，使 State/PptRecord 与实际 PPT 保持一致。
        "slides_manifest": [
            *(state.get("slides_manifest") or []),
            *[_manifest_entry_from_operation(operation) for operation in operations],
        ],
        "completed_agents": ["writer"],
        "completed_stages": ["assets"],
    }


async def skip_assets_node(state: WorkflowState) -> dict:
    """防御性跳过空的 Assets 计划。"""
    if state.get("intent") == "edit":
        return {
            "asset_apply_status": "failed",
            "workflow_error": "Edit Supervisor 选择了 Assets，但没有指定资源任务",
        }
    return {
        "asset_apply_status": "skipped",
        "completed_stages": ["assets"],
    }


async def debug_assets_join_node(_state: WorkflowState) -> dict:
    """Debug Assets 并行分支的汇合节点。"""
    return {
        "asset_apply_status": "succeeded",
        "completed_stages": ["assets"],
    }
