"""PPT 生成工具 - 根据结构化 JSON 幻灯片数据生成 PPTX 文件"""

import json
import logging

from langchain_core.tools import tool
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import settings

logger = logging.getLogger(__name__)

# 幻灯片尺寸 (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 主题配色
STYLE_THEMES = {
    "business": {
        "primary": RGBColor(0x1F, 0x3A, 0x5F),
        "accent": RGBColor(0x2E, 0x75, 0xB6),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light_text": RGBColor(0xFF, 0xFF, 0xFF),
        "title_bg": RGBColor(0x1F, 0x3A, 0x5F),
    },
    "creative": {
        "primary": RGBColor(0xE9, 0x1E, 0x63),
        "accent": RGBColor(0xFF, 0x6F, 0x00),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light_text": RGBColor(0xFF, 0xFF, 0xFF),
        "title_bg": RGBColor(0xE9, 0x1E, 0x63),
    },
    "academic": {
        "primary": RGBColor(0x1B, 0x5E, 0x20),
        "accent": RGBColor(0x43, 0xA0, 0x47),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light_text": RGBColor(0xFF, 0xFF, 0xFF),
        "title_bg": RGBColor(0x1B, 0x5E, 0x20),
    },
    "minimalist": {
        "primary": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x60, 0x60, 0x60),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light_text": RGBColor(0xFF, 0xFF, 0xFF),
        "title_bg": RGBColor(0x33, 0x33, 0x33),
    },
}

DEFAULT_THEME = {
    "primary": RGBColor(0x4F, 0x46, 0xE5),
    "accent": RGBColor(0x63, 0x66, 0xF1),
    "text": RGBColor(0x33, 0x33, 0x33),
    "light_text": RGBColor(0xFF, 0xFF, 0xFF),
    "title_bg": RGBColor(0x4F, 0x46, 0xE5),
}


@tool
async def generate_ppt(
    slides: str,
    filename: str = "presentation.pptx",
    style: str = "business",
) -> str:
    """根据结构化的幻灯片数据生成 PPTX 文件。

    slides 参数是一个 JSON 字符串，格式为幻灯片数组，每个幻灯片包含：
    - title (str): 幻灯片标题
    - bullets (list[str]): 要点列表
    - layout_hint (str): 布局类型，可选值: "title-slide", "content", "two-column"
    - speaker_notes (str, 可选): 演讲者备注

    示例:
    [
        {"title": "项目简介", "bullets": ["背景", "目标"], "layout_hint": "title-slide"},
        {"title": "技术方案", "bullets": ["方案A", "方案B"], "layout_hint": "content"}
    ]

    Args:
        slides: JSON 字符串格式的幻灯片数组
        filename: 输出文件名，默认 presentation.pptx
        style: 主题风格，可选 business/creative/academic/minimalist
    """
    try:
        # 解析 slides JSON
        slides_data = json.loads(slides)
        if not isinstance(slides_data, list) or not slides_data:
            return "错误：slides 必须是非空的 JSON 数组"

        # 确保输出目录存在
        output_dir = settings.workspace_path / "ppt_output"
        output_dir.mkdir(parents=True, exist_ok=True)

        theme = STYLE_THEMES.get(style, DEFAULT_THEME)

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        for idx, slide_data in enumerate(slides_data):
            if not isinstance(slide_data, dict):
                continue
            title = slide_data.get("title", "")
            bullets = slide_data.get("bullets", [])
            layout_hint = slide_data.get("layout_hint", "content")
            speaker_notes = slide_data.get("speaker_notes", "")

            if layout_hint == "title-slide":
                _create_title_slide(prs, title, bullets, theme)
            elif layout_hint == "two-column":
                _create_two_column_slide(prs, title, bullets, theme)
            else:
                _create_content_slide(prs, title, bullets, theme)

            # 添加演讲者备注
            if speaker_notes:
                slide = prs.slides[-1]
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = speaker_notes

        # 保存文件
        file_path = output_dir / filename
        prs.save(str(file_path))

        result = {
            "success": True,
            "file_path": str(file_path),
            "slides_count": len(slides_data),
            "message": f"PPT 生成成功: {filename}，共 {len(slides_data)} 页",
        }
        logger.info("PPT 生成成功: %s, %d 页", filename, len(slides_data))
        return json.dumps(result, ensure_ascii=False)

    except json.JSONDecodeError as e:
        return f"错误：slides JSON 解析失败: {e}"
    except (OSError, TypeError, ValueError, KeyError, AttributeError) as e:
        logger.exception("PPT 生成失败")
        return f"错误：PPT 生成失败: {e}"


# --- 内部辅助函数 ---


def _add_background(slide, color: RGBColor):
    """添加纯色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=None,
    bold=False,
    alignment=PP_ALIGN.LEFT,
):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_bullet_list(
    slide, left, top, width, height, bullets: list[str], font_size=16, color=None
):
    """添加要点列表"""
    if not bullets:
        return None
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        if color:
            p.font.color.rgb = color
        p.space_after = Pt(8)
        # 添加项目符号缩进
        p.level = 0

    return txBox


def _create_title_slide(prs, title: str, bullets: list[str], theme: dict):
    """创建封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景色
    _add_background(slide, theme["title_bg"])

    # 左侧装饰条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(2.0),
        Inches(0.15),
        Inches(3.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.fill.background()

    # 标题
    if title:
        _add_textbox(
            slide,
            Inches(1.2),
            Inches(2.2),
            Inches(10),
            Inches(1.5),
            title,
            font_size=40,
            color=theme["light_text"],
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

    # 副标题/要点
    if bullets:
        _add_bullet_list(
            slide,
            Inches(1.2),
            Inches(4.0),
            Inches(10),
            Inches(2.5),
            bullets,
            font_size=20,
            color=theme["light_text"],
        )


def _create_content_slide(prs, title: str, bullets: list[str], theme: dict):
    """创建内容页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 白色背景
    _add_background(slide, RGBColor(0xFF, 0xFF, 0xFF))

    # 顶部标题栏
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(1.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["primary"]
    shape.line.fill.background()

    # 标题文字
    if title:
        _add_textbox(
            slide,
            Inches(0.8),
            Inches(0.2),
            Inches(11),
            Inches(0.8),
            title,
            font_size=28,
            color=theme["light_text"],
            bold=True,
        )

    # 内容区域
    if bullets:
        _add_bullet_list(
            slide,
            Inches(1.0),
            Inches(1.8),
            Inches(11),
            Inches(5.0),
            bullets,
            font_size=18,
            color=theme["text"],
        )


def _create_two_column_slide(prs, title: str, bullets: list[str], theme: dict):
    """创建双栏内容页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 白色背景
    _add_background(slide, RGBColor(0xFF, 0xFF, 0xFF))

    # 顶部标题栏
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(1.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["primary"]
    shape.line.fill.background()

    # 标题文字
    if title:
        _add_textbox(
            slide,
            Inches(0.8),
            Inches(0.2),
            Inches(11),
            Inches(0.8),
            title,
            font_size=28,
            color=theme["light_text"],
            bold=True,
        )

    # 将要点分为两栏
    mid = len(bullets) // 2
    left_bullets = bullets[:mid] or bullets
    right_bullets = bullets[mid:] if mid > 0 else []

    # 左栏
    if left_bullets:
        _add_bullet_list(
            slide,
            Inches(0.8),
            Inches(1.8),
            Inches(5.5),
            Inches(5.0),
            left_bullets,
            font_size=16,
            color=theme["text"],
        )

    # 右栏
    if right_bullets:
        _add_bullet_list(
            slide,
            Inches(6.8),
            Inches(1.8),
            Inches(5.5),
            Inches(5.0),
            right_bullets,
            font_size=16,
            color=theme["text"],
        )
