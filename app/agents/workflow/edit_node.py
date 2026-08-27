"""确定性 PPT 编辑节点：汇总并行 Agent 的操作并单点写入文件。"""

import asyncio
import logging
import os
import uuid
from pathlib import Path
from typing import Any

from langchain_core.messages import HumanMessage
from pptx import Presentation

from app.agents.workflow.state import WorkflowState
from app.core.config import settings
from app.tools.chart_generator import _add_chart_slide_to_ppt
from app.tools.pexels_search import _add_image_slide
from app.tools.ppt_generator import DEFAULT_THEME, STYLE_THEMES

logger = logging.getLogger(__name__)

EDIT_NODE = "edit_node"


def _safe_path(path_value: str, *, root: Path) -> Path:
    """只允许 Workflow 操作 workspace 内的文件，阻止路径穿越。"""
    resolved_root = root.resolve()
    candidate = Path(path_value)
    if not candidate.is_absolute():
        candidate = resolved_root / candidate
    candidate = candidate.resolve()
    if not candidate.is_relative_to(resolved_root):
        raise ValueError(f"资源路径不在 workspace 内: {path_value}")
    return candidate


def _ppt_path(filename: str) -> Path:
    output_dir = (settings.workspace_path / "ppt_output").resolve()
    return _safe_path(filename, root=output_dir)


def _operation_sort_key(operation: dict[str, Any]) -> tuple:
    """并行分支返回顺序不稳定，因此提交前必须确定性排序。"""
    return (
        int(operation.get("slide_index", 10**9)),
        0 if operation.get("type") == "image" else 1,
        str(operation.get("operation_id", "")),
    )


def _apply_operation(prs: Presentation, operation: dict[str, Any]) -> None:
    operation_type = operation.get("type")
    asset_path_value = operation.get("asset_path")
    if not isinstance(asset_path_value, str) or not asset_path_value:
        raise ValueError("编辑操作缺少 asset_path")

    asset_path = _safe_path(asset_path_value, root=settings.workspace_path)
    if not asset_path.is_file():
        raise FileNotFoundError(f"编辑资源不存在: {asset_path.name}")

    style = str(operation.get("style") or "business")
    theme = STYLE_THEMES.get(style, DEFAULT_THEME)
    slide_title = str(operation.get("slide_title") or "")

    if operation_type == "image":
        _add_image_slide(prs, asset_path, slide_title, theme)
        return
    if operation_type == "chart":
        _add_chart_slide_to_ppt(prs, asset_path, slide_title, theme)
        return
    raise ValueError(f"不支持的编辑操作类型: {operation_type!r}")


def _apply_operations_atomically(
    file_path: Path,
    operations: list[dict[str, Any]],
) -> None:
    """在内存中应用全部操作，通过临时文件 + os.replace 原子提交。"""
    if not file_path.is_file():
        raise FileNotFoundError(f"PPT 文件不存在: {file_path.name}")

    temp_path = file_path.with_name(
        f".{file_path.stem}.{uuid.uuid4().hex}.tmp.pptx"
    )
    try:
        prs = Presentation(str(file_path))
        for operation in operations:
            _apply_operation(prs, operation)
        prs.save(str(temp_path))

        # 提交前重新打开，至少保证生成物是 python-pptx 可解析的 PPTX。
        Presentation(str(temp_path))
        os.replace(temp_path, file_path)
    finally:
        temp_path.unlink(missing_ok=True)


async def edit_node(state: WorkflowState) -> dict[str, Any]:
    """合并 Image/Chart 操作；本节点不调用 Agent 或 LLM。"""
    filename = state.get("filename")
    if not filename:
        raise ValueError("edit_node 缺少目标 PPT 文件名")

    operations = sorted(
        state.get("asset_operations", []),
        key=_operation_sort_key,
    )
    if not operations:
        return {
            "asset_apply_status": "skipped",
            "applied_operation_ids": [],
            "completed_agents": ["edit"],
            "messages": [
                HumanMessage(
                    content="edit_node 未收到可执行的图片或图表操作，已跳过写入",
                    name=EDIT_NODE,
                )
            ],
        }

    # 当前生产子图中 edit_node 是唯一 Writer，因此暂不加锁。未来若允许
    # 多个 Workflow 跨 Worker 修改同一 ppt_id，可在这一层接入 Redis Lock。
    await asyncio.to_thread(
        _apply_operations_atomically,
        _ppt_path(filename),
        operations,
    )

    operation_ids = [
        str(operation["operation_id"])
        for operation in operations
        if operation.get("operation_id")
    ]
    logger.info(
        "edit_node 已统一写入 PPT: filename=%s operations=%d",
        filename,
        len(operations),
    )
    return {
        "asset_apply_status": "succeeded",
        "applied_operation_ids": operation_ids,
        "completed_agents": ["edit"],
        "messages": [
            HumanMessage(
                content=f"edit_node 已统一应用 {len(operations)} 项 PPT 编辑操作",
                name=EDIT_NODE,
            )
        ],
    }
