"""图表生成工具 - 使用 matplotlib 生成数据可视化图表并插入 PPT 幻灯片"""

import json
import logging
import uuid
from pathlib import Path

import matplotlib
from langchain_core.tools import tool
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

matplotlib.use("Agg")  # 非交互式后端
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt

from app.core.config import settings
from app.tools.ppt_generator import (
    DEFAULT_THEME,
    SLIDE_WIDTH,
    STYLE_THEMES,
)

logger = logging.getLogger(__name__)

# 默认图表配色
DEFAULT_CHART_COLORS = [
    "#4F46E5",
    "#2E75B6",
    "#43A047",
    "#E91E63",
    "#FF6F00",
    "#7B1FA2",
    "#00838F",
    "#C62828",
]

# 中文字体候选列表
_CHINESE_FONT_CANDIDATES = [
    "SimHei",  # 黑体（Windows）
    "Microsoft YaHei",  # 微软雅黑（Windows）
    "WenQuanYi Micro Hei",  # 文泉驿（Linux）
    "PingFang SC",  # 苹方（macOS）
    "Noto Sans CJK SC",  # 思源（跨平台）
]

_chinese_font_initialized = False


def _setup_chinese_fonts():
    """配置 matplotlib 中文字体支持（仅初始化一次）"""
    global _chinese_font_initialized
    if _chinese_font_initialized:
        return
    _chinese_font_initialized = True

    # 获取系统所有可用字体名
    available_fonts = {f.name for f in fm.fontManager.ttflist}

    # 按优先级选择第一个可用的中文字体
    selected = None
    for candidate in _CHINESE_FONT_CANDIDATES:
        if candidate in available_fonts:
            selected = candidate
            break

    if selected:
        plt.rcParams["font.sans-serif"] = [selected, "DejaVu Sans"]
        logger.info("matplotlib 中文字体设置为: %s", selected)
    else:
        # 回退方案：使用系统默认 sans-serif
        plt.rcParams["font.sans-serif"] = ["sans-serif"]
        logger.warning("未找到中文字体，图表中文可能显示异常")

    plt.rcParams["axes.unicode_minus"] = False  # 正确显示负号


def _hex_to_tuple(hex_color: str) -> tuple:
    """将十六进制颜色字符串转为 RGB 元组"""
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) / 255.0 for i in (0, 2, 4))


def _get_chart_colors(config: dict) -> list:
    """获取图表颜色列表"""
    custom_colors = config.get("colors", [])
    if custom_colors:
        return [_hex_to_tuple(c) for c in custom_colors]
    return [_hex_to_tuple(c) for c in DEFAULT_CHART_COLORS]


def _generate_bar_chart(data: dict, output_path: Path, colors: list, title: str):
    """生成柱状图"""
    fig, ax = plt.subplots(figsize=(10, 5.5))
    labels = data.get("labels", [])
    values = data.get("values", [])

    bars = ax.bar(
        range(len(labels)),
        values,
        color=colors[: len(labels)],
        edgecolor="white",
        linewidth=0.5,
        width=0.6,
    )

    # 在柱子上方显示数值
    for bar, val in zip(bars, values):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height() + max(values) * 0.02,
            str(val),
            ha="center",
            va="bottom",
            fontsize=11,
            fontweight="bold",
        )

    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_title(title, fontsize=14, fontweight="bold", pad=15)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_alpha(0.3)
    ax.spines["bottom"].set_alpha(0.3)
    ax.yaxis.grid(True, alpha=0.2)
    ax.set_axisbelow(True)

    plt.tight_layout()
    fig.savefig(
        str(output_path),
        dpi=200,
        bbox_inches="tight",
        facecolor="white",
        edgecolor="none",
    )
    plt.close(fig)


def _generate_pie_chart(data: dict, output_path: Path, colors: list, title: str):
    """生成饼图"""
    fig, ax = plt.subplots(figsize=(8, 6))
    labels = data.get("labels", [])
    values = data.get("values", [])

    wedges, texts, autotexts = ax.pie(
        values,
        labels=labels,
        colors=colors[: len(labels)],
        autopct="%1.1f%%",
        startangle=90,
        pctdistance=0.75,
        textprops={"fontsize": 11},
    )
    for autotext in autotexts:
        autotext.set_fontsize(10)
        autotext.set_fontweight("bold")

    ax.set_title(title, fontsize=14, fontweight="bold", pad=15)

    plt.tight_layout()
    fig.savefig(
        str(output_path),
        dpi=200,
        bbox_inches="tight",
        facecolor="white",
        edgecolor="none",
    )
    plt.close(fig)


def _generate_line_chart(data: dict, output_path: Path, colors: list, title: str):
    """生成折线图"""
    fig, ax = plt.subplots(figsize=(10, 5.5))
    labels = data.get("labels", [])
    values = data.get("values", [])

    ax.plot(
        range(len(labels)),
        values,
        color=colors[0],
        marker="o",
        markersize=8,
        linewidth=2.5,
        markerfacecolor="white",
        markeredgewidth=2,
        markeredgecolor=colors[0],
    )

    # 在数据点上显示数值
    for i, val in enumerate(values):
        ax.annotate(
            str(val),
            (i, val),
            textcoords="offset points",
            xytext=(0, 12),
            ha="center",
            fontsize=10,
            fontweight="bold",
        )

    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_title(title, fontsize=14, fontweight="bold", pad=15)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_alpha(0.3)
    ax.spines["bottom"].set_alpha(0.3)
    ax.yaxis.grid(True, alpha=0.2)
    ax.set_axisbelow(True)

    # 填充折线下方区域（浅色）
    ax.fill_between(range(len(labels)), values, alpha=0.1, color=colors[0])

    plt.tight_layout()
    fig.savefig(
        str(output_path),
        dpi=200,
        bbox_inches="tight",
        facecolor="white",
        edgecolor="none",
    )
    plt.close(fig)


def _generate_area_chart(data: dict, output_path: Path, colors: list, title: str):
    """生成面积图"""
    fig, ax = plt.subplots(figsize=(10, 5.5))
    labels = data.get("labels", [])
    values = data.get("values", [])

    ax.fill_between(range(len(labels)), values, alpha=0.3, color=colors[0])
    ax.plot(
        range(len(labels)),
        values,
        color=colors[0],
        linewidth=2.5,
        marker="o",
        markersize=6,
    )

    # 在数据点上显示数值
    for i, val in enumerate(values):
        ax.annotate(
            str(val),
            (i, val),
            textcoords="offset points",
            xytext=(0, 12),
            ha="center",
            fontsize=10,
            fontweight="bold",
        )

    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_title(title, fontsize=14, fontweight="bold", pad=15)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_alpha(0.3)
    ax.spines["bottom"].set_alpha(0.3)
    ax.yaxis.grid(True, alpha=0.2)
    ax.set_axisbelow(True)

    plt.tight_layout()
    fig.savefig(
        str(output_path),
        dpi=200,
        bbox_inches="tight",
        facecolor="white",
        edgecolor="none",
    )
    plt.close(fig)


def _add_chart_slide_to_ppt(
    prs: Presentation, image_path: Path, slide_title: str, theme: dict
):
    """向 PPT 中添加一张包含图表图片的幻灯片"""
    from pptx.enum.shapes import MSO_SHAPE

    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 白色背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 顶部标题栏
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(1.0),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["primary"]
    shape.line.fill.background()

    # 标题文字
    if slide_title:
        txBox = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(0.15),
            Inches(11),
            Inches(0.7),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(26)
        p.font.color.rgb = theme["light_text"]
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    # 插入图表图片（居中）
    chart_left = Inches(1.5)
    chart_top = Inches(1.3)
    chart_width = Inches(10.0)
    chart_height = Inches(5.5)
    slide.shapes.add_picture(
        str(image_path), chart_left, chart_top, chart_width, chart_height
    )


def _validate_chart_config(config: object) -> dict:
    """校验模型生成的图表配置，避免错误结构进入 matplotlib。"""
    if not isinstance(config, dict):
        raise ValueError(
            "chart_config 顶层必须是 JSON 对象，不能是数组；多张图表请分别调用工具"
        )

    chart_type = config.get("chart_type", "bar")
    if chart_type not in {"bar", "pie", "line", "area"}:
        raise ValueError(
            f"不支持的图表类型 {chart_type!r}，可选: bar/pie/line/area"
        )

    data = config.get("data")
    if not isinstance(data, dict):
        raise ValueError(
            "chart_config.data 必须是包含 labels 和 values 的 JSON 对象，不能是数组"
        )

    labels = data.get("labels")
    values = data.get("values")
    if not isinstance(labels, list) or not labels:
        raise ValueError("chart_config.data.labels 必须是非空数组")
    if not isinstance(values, list) or not values:
        raise ValueError("chart_config.data.values 必须是非空数组")
    if len(labels) != len(values):
        raise ValueError("chart_config.data.labels 与 values 的数量必须一致")
    if not all(
        isinstance(value, (int, float)) and not isinstance(value, bool)
        for value in values
    ):
        raise ValueError("chart_config.data.values 只能包含数字")

    return config


def _prepare_chart_asset(config: dict) -> tuple[Path, str, str, str]:
    """渲染图表图片，但不打开或修改任何 PPT 文件。"""
    chart_type = config.get("chart_type", "bar")
    data = config.get("data", {})
    title = config.get("title", "")
    slide_title = config.get("slide_title", title)

    _setup_chinese_fonts()
    output_dir = settings.workspace_path / "ppt_output"
    charts_dir = output_dir / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)
    chart_image = charts_dir / f"chart_{uuid.uuid4().hex[:8]}.png"
    colors = _get_chart_colors(config)

    generators = {
        "bar": _generate_bar_chart,
        "pie": _generate_pie_chart,
        "line": _generate_line_chart,
        "area": _generate_area_chart,
    }
    generator = generators.get(chart_type)
    if not generator:
        raise ValueError(
            f"不支持的图表类型 '{chart_type}'，可选: bar/pie/line/area"
        )

    generator(data, chart_image, colors, title)
    return chart_image, chart_type, title, slide_title


@tool
async def prepare_chart_operation(filename: str, chart_config: str) -> str:
    """准备一项 PPT 图表操作，但不直接修改 PPT 文件。

    返回的 operation 会由 Workflow 的确定性 ppt_writer_node 汇总并统一写入。

    Args:
        filename: 目标 PPT 文件名，仅用于标识操作目标
        chart_config: 单个 JSON 对象的字符串，data 必须是包含 labels 和 values
            的对象。多张图表需要分别调用本工具，不能传入 JSON 数组。
    """
    try:
        config = _validate_chart_config(json.loads(chart_config))
        output_dir = settings.workspace_path / "ppt_output"
        if not (output_dir / filename).exists():
            return json.dumps(
                {"success": False, "message": f"PPT 文件不存在: {filename}"},
                ensure_ascii=False,
            )

        chart_image, chart_type, title, slide_title = _prepare_chart_asset(config)
        operation = {
            "operation_id": f"chart-{uuid.uuid4().hex}",
            "type": "chart",
            "filename": filename,
            "asset_path": str(chart_image),
            "slide_title": slide_title,
            "style": config.get("style", "business"),
            "chart_type": chart_type,
            "chart_title": title,
        }
        return json.dumps(
            {
                "success": True,
                "operation": operation,
                "message": f"图表资源已准备，等待统一写入: {chart_type}",
            },
            ensure_ascii=False,
        )
    except json.JSONDecodeError as exc:
        return json.dumps(
            {"success": False, "message": f"chart_config JSON 解析失败: {exc}"},
            ensure_ascii=False,
        )
    except ValueError as exc:
        return json.dumps(
            {"success": False, "message": f"chart_config 格式错误: {exc}"},
            ensure_ascii=False,
        )
    except Exception as exc:
        logger.exception("准备图表操作失败")
        return json.dumps(
            {"success": False, "message": f"准备图表操作失败: {exc}"},
            ensure_ascii=False,
        )


@tool
async def add_chart_slide(
    filename: str,
    chart_config: str,
) -> str:
    """生成数据可视化图表并插入为新的PPT幻灯片。

    chart_config JSON格式示例：
    {
        "chart_type": "bar",
        "title": "图表标题",
        "data": {
            "labels": ["Q1", "Q2", "Q3", "Q4"],
            "values": [100, 200, 150, 300]
        },
        "colors": ["#4F46E5", "#2E75B6"],
        "slide_title": "幻灯片标题",
        "style": "business"
    }

    支持的图表类型（chart_type）：
    - "bar": 柱状图
    - "pie": 饼图
    - "line": 折线图
    - "area": 面积图

    Args:
        filename: PPT文件名（在ppt_output目录下）
        chart_config: JSON字符串，描述图表配置
    """
    try:
        config = _validate_chart_config(json.loads(chart_config))
        chart_type = config.get("chart_type", "bar")
        style = config.get("style", "business")
        theme = STYLE_THEMES.get(style, DEFAULT_THEME)

        # 解析文件路径
        output_dir = settings.workspace_path / "ppt_output"
        file_path = output_dir / filename
        if not file_path.exists():
            return f"错误：PPT 文件不存在: {filename}"

        chart_image, chart_type, _title, slide_title = _prepare_chart_asset(config)

        # 打开 PPT，插入图表幻灯片
        prs = Presentation(str(file_path))
        _add_chart_slide_to_ppt(prs, chart_image, slide_title, theme)
        prs.save(str(file_path))

        result = {
            "success": True,
            "chart_type": chart_type,
            "chart_image": str(chart_image),
            "message": f"图表幻灯片已添加到 {filename}（类型: {chart_type}）",
        }
        logger.info("图表幻灯片添加成功: %s, 类型: %s", filename, chart_type)
        return json.dumps(result, ensure_ascii=False)

    except json.JSONDecodeError as e:
        return f"错误：chart_config JSON 解析失败: {e}"
    except ValueError as e:
        return f"错误：chart_config 格式错误: {e}"
    except Exception as e:
        logger.error("图表生成失败: %s", e)
        return f"错误：图表生成失败: {e}"
