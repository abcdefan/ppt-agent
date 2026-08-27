"""PPT 增强美化"""

import json
import logging
from pathlib import Path

from langchain_core.tools import tool
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from app.core.config import settings
from app.tools.ppt_generator import (
    DEFAULT_THEME,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    STYLE_THEMES,
)

logger = logging.getLogger(__name__)


# --- 通用辅助 ---


def _get_theme(style: str) -> dict:
    return STYLE_THEMES.get(style, DEFAULT_THEME)


def _resolve_ppt_path(filename: str) -> Path:
    return settings.workspace_path / "ppt_output" / filename


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=None,
    bold=False,
    alignment=PP_ALIGN.LEFT,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_header_bar(slide, title: str, theme: dict):
    """添加统一的顶部标题栏"""
    header_h = Inches(1.1)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        header_h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["primary"]
    shape.line.fill.background()
    if title:
        # 设置标题文字垂直居中
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(11), header_h - Inches(0.3)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.color.rgb = theme["light_text"]
        p.font.bold = True


# ============================================================
# action="layout" — 高级布局（timeline / comparison / stats / card-grid）
# ============================================================


def _create_timeline(prs, config: dict, theme: dict):
    """时间线布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _add_header_bar(slide, config.get("title", ""), theme)

    items = config.get("items", [])
    if not items:
        return

    n = len(items)
    line_y = Inches(3.8)
    start_x, end_x = Inches(1.5), Inches(11.8)
    spacing = (end_x - start_x) / max(n - 1, 1) if n > 1 else 0

    # 连接线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        start_x,
        line_y,
        end_x - start_x,
        Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme["accent"]
    line.line.fill.background()

    for i, item in enumerate(items):
        x = start_x + spacing * i if n > 1 else (start_x + end_x) // 2

        # 节点圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x - Inches(0.2),
            line_y - Inches(0.15),
            Inches(0.4),
            Inches(0.4),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme["primary"]
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 时间标签（上方）
        time_text = item.get("time", "")
        if time_text:
            _add_textbox(
                slide,
                x - Inches(0.6),
                line_y - Inches(0.9),
                Inches(1.2),
                Inches(0.6),
                time_text,
                font_size=12,
                color=theme["primary"],
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )
        # 描述文字（下方）
        desc = item.get("text", "")
        if desc:
            _add_textbox(
                slide,
                x - Inches(0.7),
                line_y + Inches(0.5),
                Inches(1.4),
                Inches(1.5),
                desc,
                font_size=10,
                color=theme["text"],
                alignment=PP_ALIGN.CENTER,
            )


def _create_comparison(prs, config: dict, theme: dict):
    """对比布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _add_header_bar(slide, config.get("title", ""), theme)

    left_header = config.get("left_header", "")
    left_items = config.get("left_items", [])
    right_header = config.get("right_header", "")
    right_items = config.get("right_items", [])

    # 分隔线
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.55),
        Inches(1.3),
        Pt(2),
        Inches(5.5),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = theme["accent"]
    divider.line.fill.background()

    for side, header, items, x_pos in [
        ("left", left_header, left_items, Inches(0.8)),
        ("right", right_header, right_items, Inches(7.0)),
    ]:
        # 标题背景条
        hbg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x_pos,
            Inches(1.3),
            Inches(5.3),
            Inches(0.55),
        )
        hbg.fill.solid()
        hbg.fill.fore_color.rgb = theme["accent"]
        hbg.line.fill.background()
        if header:
            tf = hbg.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

        # 内容列表
        if items:
            txBox = slide.shapes.add_textbox(
                x_pos + Inches(0.2), Inches(2.1), Inches(4.8), Inches(4.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"•  {item}"
                p.font.size = Pt(13)
                p.font.color.rgb = theme["text"]
                p.space_before = Pt(4)
                p.space_after = Pt(6)


def _create_stats(prs, config: dict, theme: dict):
    """数据展示布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _add_header_bar(slide, config.get("title", ""), theme)

    stats = config.get("stats", [])
    if not stats:
        return

    n = len(stats)
    cols = min(n, 4)
    col_w = (Inches(10.5)) // cols
    start_x = Inches(1.5)

    for i, stat in enumerate(stats):
        col = i % cols
        x = start_x + col_w * col + Inches(0.2)
        y = Inches(1.5)

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            col_w - Inches(0.4),
            Inches(4.5),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        card.line.fill.background()

        # 大数字
        number = stat.get("number", "")
        if number:
            _add_textbox(
                slide,
                x,
                y + Inches(0.8),
                col_w - Inches(0.4),
                Inches(1.5),
                number,
                font_size=42,
                color=theme["primary"],
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )
        # 标签
        label = stat.get("label", "")
        if label:
            _add_textbox(
                slide,
                x,
                y + Inches(2.5),
                col_w - Inches(0.4),
                Inches(0.8),
                label,
                font_size=14,
                color=theme["text"],
                alignment=PP_ALIGN.CENTER,
            )
        # 底部装饰线
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.5),
            y + Inches(3.6),
            col_w - Inches(1.4),
            Pt(3),
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = theme["accent"]
        accent_line.line.fill.background()


def _create_card_grid(prs, config: dict, theme: dict):
    """卡片网格布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _add_header_bar(slide, config.get("title", ""), theme)

    cards = config.get("cards", [])
    if not cards:
        return

    cols = min(len(cards), 3)
    card_w, card_h = Inches(3.6), Inches(2.5)
    gap_x = Inches(0.5)
    total_w = card_w * cols + gap_x * (cols - 1)
    start_x = (SLIDE_WIDTH - total_w) // 2
    start_y = Inches(1.5)

    colors = [theme["primary"], theme["accent"]]
    _p = str(theme["primary"])
    colors.append(
        RGBColor(
            min(int(_p[0:2], 16) + 40, 255),
            min(int(_p[2:4], 16) + 40, 255),
            min(int(_p[4:6], 16) + 40, 255),
        )
    )

    for i, card in enumerate(cards):
        col, row = i % cols, i // cols
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + Inches(0.4)) * row

        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            card_w,
            card_h,
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = colors[i % len(colors)]
        card_shape.line.fill.background()

        card_title = card.get("title", "")
        if card_title:
            _add_textbox(
                slide,
                x + Inches(0.3),
                y + Inches(0.5),
                card_w - Inches(0.6),
                Inches(0.6),
                card_title,
                font_size=16,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                bold=True,
            )
        card_text = card.get("text", "")
        if card_text:
            _add_textbox(
                slide,
                x + Inches(0.3),
                y + Inches(1.2),
                card_w - Inches(0.6),
                Inches(1.0),
                card_text,
                font_size=11,
                color=RGBColor(0xFF, 0xFF, 0xFF),
            )


_LAYOUT_CREATORS = {
    "timeline": _create_timeline,
    "comparison": _create_comparison,
    "stats": _create_stats,
    "card-grid": _create_card_grid,
}


# ============================================================
# action="decorate" — 装饰元素（step_indicators / process_flow）
# ============================================================


def _add_step_indicators(slide, steps: list, active_step: int, theme: dict):
    """步骤指示器（编号圆圈 + 连接线，位于底部）"""
    if not steps:
        return
    n = len(steps)
    circle_size = Inches(0.5)
    spacing = Inches(1.8)
    total_w = spacing * (n - 1) + circle_size
    start_x = (SLIDE_WIDTH - total_w) // 2
    start_y = Inches(6.3)

    for i, step_text in enumerate(steps):
        x = start_x + spacing * i

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x,
            start_y,
            circle_size,
            circle_size,
        )
        if i == active_step:
            circle.fill.solid()
            circle.fill.fore_color.rgb = theme["accent"]
        elif i < active_step:
            circle.fill.solid()
            circle.fill.fore_color.rgb = theme["primary"]
        else:
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
        circle.line.fill.background()

        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 步骤文字
        _add_textbox(
            slide,
            x - Inches(0.5),
            start_y + circle_size + Inches(0.05),
            Inches(1.5),
            Inches(0.35),
            step_text,
            font_size=9,
            color=theme["text"],
            alignment=PP_ALIGN.CENTER,
        )

        # 连接线
        if i < n - 1:
            ln = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + circle_size,
                start_y + circle_size // 2 - Pt(1.5),
                spacing - circle_size,
                Pt(3),
            )
            ln.fill.solid()
            ln.fill.fore_color.rgb = (
                theme["primary"] if i < active_step else RGBColor(0xD0, 0xD0, 0xD0)
            )
            ln.line.fill.background()


def _add_process_flow(slide, steps: list, theme: dict):
    """流程图（圆角矩形 + 箭头，位于中央）"""
    if not steps:
        return
    n = len(steps)
    box_w, box_h = Inches(1.6), Inches(0.7)
    arrow_w = Inches(0.5)
    total_w = box_w * n + arrow_w * (n - 1)
    start_x = (SLIDE_WIDTH - total_w) // 2
    start_y = Inches(3.2)

    for i, step_text in enumerate(steps):
        x = start_x + (box_w + arrow_w) * i

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            start_y,
            box_w,
            box_h,
        )
        box.fill.solid()
        box.fill.fore_color.rgb = theme["primary"]
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_text
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if i < n - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + Inches(0.05),
                start_y + box_h // 2 - Inches(0.12),
                arrow_w - Inches(0.1),
                Inches(0.24),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = theme["accent"]
            arrow.line.fill.background()


# ============================================================
# action="beautify" — 视觉美化（渐变背景 + 页码 + 装饰 + 排版）
# ============================================================


def _apply_gradient_bg(slide, c1_hex: str, c2_hex: str):
    """通过 OOXML XML 设置线性渐变背景"""
    bg_elem = slide.background._element
    for child in list(bg_elem):
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
        if tag in ("bgPr", "bgRef"):
            bg_elem.remove(child)

    bgPr = etree.SubElement(bg_elem, qn("p:bgPr"))
    gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

    gs1 = etree.SubElement(gsLst, qn("a:gs"), attrib={"pos": "0"})
    etree.SubElement(gs1, qn("a:srgbClr"), attrib={"val": c1_hex})

    gs2 = etree.SubElement(gsLst, qn("a:gs"), attrib={"pos": "100000"})
    etree.SubElement(gs2, qn("a:srgbClr"), attrib={"val": c2_hex})

    etree.SubElement(
        gradFill, qn("a:lin"), attrib={"ang": str(270 * 60000), "scaled": "1"}
    )
    etree.SubElement(bgPr, qn("a:effectLst"))


def _add_page_number(slide, num: int, total: int):
    """右下角页码"""
    txBox = slide.shapes.add_textbox(
        Inches(11.0), Inches(7.0), Inches(2.0), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.RIGHT


def _add_corner_accents(slide, theme: dict):
    """角落装饰圆"""
    for x, y, w, alpha in [
        (SLIDE_WIDTH - Inches(1.5), Inches(-0.5), Inches(2.0), 15000),
        (Inches(-0.8), SLIDE_HEIGHT - Inches(1.2), Inches(1.8), 10000),
    ]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, w)
        c.fill.solid()
        c.fill.fore_color.rgb = theme["accent"]
        c.line.fill.background()
        # 设置透明度
        sp = c._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is not None:
            sf = spPr.find(qn("a:solidFill"))
            if sf is not None:
                sc = sf.find(qn("a:srgbClr"))
                if sc is not None:
                    existing = sc.find(qn("a:alpha"))
                    if existing is not None:
                        sc.remove(existing)
                    etree.SubElement(sc, qn("a:alpha"), attrib={"val": str(alpha)})


def _is_title_slide(slide) -> bool:
    """判断是否为封面页"""
    bg_elem = slide.background._element
    bgPr = bg_elem.find(qn("p:bgPr"))
    if bgPr is not None:
        sf = bgPr.find(qn("a:solidFill"))
        if sf is not None:
            sc = sf.find(qn("a:srgbClr"))
            if sc is not None and sc.get("val", "") != "FFFFFF":
                return True
        if bgPr.find(qn("a:gradFill")) is not None:
            return True
    return False


# ============================================================
# 主工具入口
# ============================================================


@tool
async def enhance_ppt(
    filename: str,
    action: str,
    options: str = "{}",
) -> str:
    """对PPT进行增强操作：添加高级布局、装饰元素或视觉美化。

    action 可选值：

    1. "layout" — 添加高级布局幻灯片
       options.type 可选: "timeline", "comparison", "stats", "card-grid"
       示例: {"type": "timeline", "title": "...", "items": [{"time": "2024", "text": "..."}]}
       示例: {"type": "comparison", "title": "...", "left_header": "A", "left_items": [...], "right_header": "B", "right_items": [...]}
       示例: {"type": "stats", "title": "...", "stats": [{"number": "99%", "label": "满意度"}]}
       示例: {"type": "card-grid", "title": "...", "cards": [{"title": "...", "text": "..."}]}

    2. "decorate" — 添加装饰元素到指定幻灯片
       options 格式: {"slide_index": 0, "type": "step_indicators|process_flow", ...}
       step_indicators: {"slide_index": 0, "type": "step_indicators", "steps": ["a","b","c"], "active_step": 1}
       process_flow: {"slide_index": 0, "type": "process_flow", "steps": ["收集","分析","执行"]}

    3. "beautify" — 视觉美化（渐变背景 + 角落装饰 + 页码 + 排版优化）
       options 格式: {"style": "business"}  (style 可选，默认 business)

    Args:
        filename: PPT文件名（在ppt_output目录下）
        action: 操作类型 "layout" / "decorate" / "beautify"
        options: JSON字符串，各 action 的配置参数
    """
    try:
        # 兼容 ainvoke（可能传入 dict）和直接调用（传入 JSON str）
        if isinstance(options, dict):
            cfg = options
        elif options:
            cfg = json.loads(options)
        else:
            cfg = {}
        style = cfg.get("style", "business")
        theme = _get_theme(style)

        file_path = _resolve_ppt_path(filename)
        if not file_path.exists():
            return f"错误：PPT 文件不存在: {filename}"

        prs = Presentation(str(file_path))

        # --- layout ---
        if action == "layout":
            layout_type = cfg.get("type", "")
            creator = _LAYOUT_CREATORS.get(layout_type)
            if not creator:
                avail = ", ".join(_LAYOUT_CREATORS.keys())
                return f"错误：不支持的布局 '{layout_type}'，可选: {avail}"
            creator(prs, cfg, theme)
            prs.save(str(file_path))
            return json.dumps(
                {"success": True, "message": f"布局已添加: {layout_type}"},
                ensure_ascii=False,
            )

        # --- decorate ---
        elif action == "decorate":
            deco_type = cfg.get("type", "")
            slide_index = cfg.get("slide_index", 0)
            total = len(prs.slides)
            if slide_index < 0 or slide_index >= total:
                return f"错误：slide_index {slide_index} 超出范围（共 {total} 页）"
            slide = prs.slides[slide_index]

            if deco_type == "step_indicators":
                _add_step_indicators(
                    slide, cfg.get("steps", []), cfg.get("active_step", 0), theme
                )
            elif deco_type == "process_flow":
                _add_process_flow(slide, cfg.get("steps", []), theme)
            else:
                return f"错误：不支持的装饰类型 '{deco_type}'，可选: step_indicators, process_flow"

            prs.save(str(file_path))
            return json.dumps(
                {"success": True, "message": f"装饰已添加: {deco_type}"},
                ensure_ascii=False,
            )

        # --- beautify ---
        elif action == "beautify":
            total = len(prs.slides)
            _p = str(theme["primary"])
            pr_val, pg_val, pb_val = (
                int(_p[0:2], 16),
                int(_p[2:4], 16),
                int(_p[4:6], 16),
            )
            light_hex = f"{min(pr_val + 180, 255):02X}{min(pg_val + 180, 255):02X}{min(pb_val + 180, 255):02X}"

            for idx, slide in enumerate(prs.slides):
                is_title = _is_title_slide(slide)

                # 内容页：浅色渐变背景
                if not is_title:
                    _apply_gradient_bg(slide, "FFFFFF", light_hex)
                    _add_corner_accents(slide, theme)

                # 页码（跳过封面）
                if not is_title and idx > 0:
                    _add_page_number(slide, idx, total - 1)

            prs.save(str(file_path))
            return json.dumps(
                {"success": True, "message": f"美化完成: {total} 页"},
                ensure_ascii=False,
            )

        else:
            return f"错误：不支持的 action '{action}'，可选: layout, decorate, beautify"

    except json.JSONDecodeError as e:
        return f"错误：options JSON 解析失败: {e}"
    except Exception as e:
        logger.error("enhance_ppt 失败: %s", e)
        return f"错误：操作失败: {e}"
